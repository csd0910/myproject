import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()
# Set 16:9 layout
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Title Slide ---
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "長期欠品フラグ更新業務\nDX（Python自動化）提案書"
subtitle.text = "手作業とマニュアルの完全解析に基づく抜本的改善アプローチ"

def add_comparative_slide(title_text, as_is_texts, to_be_texts):
    layout = prs.slide_layouts[5] # Blank with title
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title_text
    
    # Left Box (As-Is)
    left_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.5))
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = RGBColor(240, 200, 200) # Light red
    left_shape.line.color.rgb = RGBColor(200, 50, 50)
    
    tf = left_shape.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "【現在の手作業（As-Is）】"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(150, 0, 0)
    
    for text in as_is_texts:
        p = tf.add_paragraph()
        p.text = "・" + text
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(10)

    # Right Box (To-Be)
    right_shape = slide.shapes.add_shape(1, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5))
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = RGBColor(200, 240, 200) # Light green
    right_shape.line.color.rgb = RGBColor(50, 200, 50)
    
    tf = right_shape.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "【Python自動化後（To-Be）】"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 100, 0)
    
    for text in to_be_texts:
        p = tf.add_paragraph()
        p.text = "・" + text
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(10)
        
    # Arrow in middle
    arrow = slide.shapes.add_shape(33, Inches(6.5), Inches(4.0), Inches(0.3), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
    arrow.line.fill.background()

# --- Slide 2: Data Update ---
add_comparative_slide(
    "業務①：CSVデータの流し込みと数式展開（データ更新）",
    [
        "システムから抽出した『品目』『未入荷』のCSVデータを、人間が手作業でシートにコピペしている。",
        "『①長期欠品』シートのI列～AF列（20列以上）のVLOOKUP数式を、数十万行下までオートフィルでコピー（約530万セルの更新）。",
        "Excelが処理限界を超え、「1分ほど固まる」というマニュアル記載があるほど重い。",
        "フリーズを防ぐため、さらに「数式を値として貼り付け直す」という無駄な後処理が発生。"
    ],
    [
        "Python (pandas) が2つのCSVデータをバックグラウンドで瞬時に読み込む。",
        "ExcelのVLOOKUP関数を使わず、Pythonのメモリ上で直接データを結合（マージ）する。",
        "「1分固まる」処理が【数秒～10秒程度】で完了し、フリーズは一切起こらない。",
        "計算式ではなく「値」だけを持った完成版のExcelファイルを一発で出力するため、値貼り付け作業が消滅。"
    ]
)

# --- Slide 3: Flag Logic ---
add_comparative_slide(
    "業務②：欠品フラグ解除の判定・リスト化",
    [
        "人間がCtrl+Fを使って商品を目視検索したり、AC～AE列の『解除OK』という文字をオートフィルタ等で目視で探している。",
        "マニュアルには「発注ランクが9997なら解除」「有効在庫>0かつ納期未定なら解除」といった明確なIF条件があるが、Excelの数式に依存している。",
        "解除対象の商品コードと商品名をコピーし、手動で『②長期欠品解除』シートに貼り付けて移動させている。",
        "目視確認と手作業のコピペが介在するため、見落としや貼り付けミス（ヒューマンエラー）のリスクが非常に高い。"
    ],
    [
        "マニュアルのIF条件（ランク9997なら終売等）をPythonのプログラム内にルールとして完全移植。",
        "データの結合時に、Pythonが数万行の全データに対して瞬時にルール判定を実行。",
        "「解除OK」となった対象データだけを自動で抽出し、『②長期欠品解除』シートに最初から書き出しておく。",
        "担当者は【Pythonが出力した結果（②長期欠品解除シート）を最終確認するだけ】となり、検索とコピペの手作業がゼロになる。"
    ]
)

out_path = r"C:\Users\フォーレスト026\MyProject\業務自動化分析ツール\AutoAnalysisLogs\daily_reports\Python自動化提案_AsIs_ToBe.pptx"
prs.save(out_path)
print(f"Presentation saved to {out_path}")

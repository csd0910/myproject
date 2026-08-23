import os
import re
import collections
import collections.abc
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def parse_mermaid_to_flowchart(slide, mermaid_text, y_start=1.5):
    # 簡単なMermaid解析（A[text] --> B["text"]）
    nodes = {}
    edges = []
    
    for line in mermaid_text.split('\n'):
        line = line.strip()
        if '-->' in line:
            parts = line.split('-->')
            if len(parts) == 2:
                left = parts[0].strip()
                right = parts[1].strip()
                
                # left node
                l_id_match = re.match(r'([A-Za-z0-9]+)(?:\[|"\b)(.+?)(?:\]|"\b)?$', left.replace('["', '[').replace('"]', ']'))
                if l_id_match:
                    l_id = l_id_match.group(1)
                    l_text = l_id_match.group(2).strip('"')
                else:
                    l_id = left; l_text = left
                nodes[l_id] = l_text
                
                # right node
                r_id_match = re.match(r'([A-Za-z0-9]+)(?:\[|"\b)(.+?)(?:\]|"\b)?$', right.replace('["', '[').replace('"]', ']'))
                if r_id_match:
                    r_id = r_id_match.group(1)
                    r_text = r_id_match.group(2).strip('"')
                else:
                    r_id = right; r_text = right
                nodes[r_id] = r_text
                
                edges.append((l_id, r_id))
    
    # 描画
    if not edges:
        tb = slide.shapes.add_textbox(Inches(1), Inches(y_start), Inches(11), Inches(5))
        tb.text_frame.text = mermaid_text
        return

    # ノードをシーケンシャルに並べる
    seq = []
    current = edges[0][0]
    while True:
        if current not in seq:
            seq.append(current)
        next_nodes = [e[1] for e in edges if e[0] == current]
        if next_nodes:
            current = next_nodes[0]
        else:
            break
            
    x_offset = 0.5
    for s_id in seq:
        text = nodes.get(s_id, s_id)
        # 四角形
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_offset), Inches(y_start), Inches(1.8), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(220, 240, 255)
        shape.line.color.rgb = RGBColor(0, 0, 0)
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.name = "Meiryo UI"
        p.alignment = PP_ALIGN.CENTER
        
        x_offset += 2.0
        if x_offset > 11:
            x_offset = 0.5
            y_start += 1.5

def extract_steps(soup):
    steps = []
    text = soup.get_text(separator='\n', strip=True)
    step_matches = re.split(r'(Step \d+:.*?)\n', text)
    
    current_step = None
    for part in step_matches:
        if part.startswith('Step'):
            if current_step:
                steps.append(current_step)
            current_step = {"title": part, "content": ""}
        elif current_step:
            current_step["content"] += part + "\n"
    if current_step:
        steps.append(current_step)
    return steps

def generate():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    html_files = [
        r"C:\AutoAnalysisLogs\daily_reports\dx_analysis_report_20260810.html",
        r"C:\AutoAnalysisLogs\daily_reports\dx_analysis_report_20260811.html"
    ]
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tb.text_frame.text = "DX分析レポート 完全統合版"
    tb.text_frame.paragraphs[0].font.size = Pt(54)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.name = "Meiryo UI"
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    for html_file in html_files:
        if not os.path.exists(html_file):
            continue
        date_str = os.path.basename(html_file).replace('dx_analysis_report_', '').replace('.html', '')
        
        with open(html_file, 'r', encoding='utf-8') as f:
            html = f.read()
        soup = BeautifulSoup(html, 'html.parser')
        raw_text = soup.get_text(separator='\n', strip=True)
        
        # 1. サマリー
        slide_summary = prs.slides.add_slide(prs.slide_layouts[5])
        slide_summary.shapes.title.text = f"[{date_str}] 分析サマリー＆使用ファイル"
        slide_summary.shapes.title.text_frame.paragraphs[0].font.name = "Meiryo UI"
        
        tb_sum = slide_summary.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
        tf_sum = tb_sum.text_frame
        
        files_match = re.search(r'操作対象ファイル/システム:(.*?)手作業の実行時間', raw_text, re.DOTALL)
        if files_match:
            p = tf_sum.add_paragraph()
            p.text = "【使用した主要ファイル】\n" + files_match.group(1).strip()
            p.font.name = "Meiryo UI"
            p.font.size = Pt(16)
        
        time_match = re.search(r'手作業の実行時間.*?(\d+秒).*?想定削減時間.*?(\d+秒)', raw_text, re.DOTALL)
        if time_match:
            p = tf_sum.add_paragraph()
            p.text = f"\n【時間コスト分析】\n・現在の手作業: {time_match.group(1)}\n・自動化後の想定: {time_match.group(2)}"
            p.font.name = "Meiryo UI"
            p.font.size = Pt(20)
            p.font.bold = True

        # 2. フロー図（現在の処理フロー）
        slide_flow_cur = prs.slides.add_slide(prs.slide_layouts[5])
        slide_flow_cur.shapes.title.text = f"[{date_str}] 現在の処理フロー図"
        slide_flow_cur.shapes.title.text_frame.paragraphs[0].font.name = "Meiryo UI"
        
        flow_cur_match = re.search(r'現在の処理フロー図.*?graph TD(.*?)自動化後のフロー図', raw_text, re.DOTALL)
        if flow_cur_match:
            parse_mermaid_to_flowchart(slide_flow_cur, flow_cur_match.group(1))

        # 3. フロー図（自動化後）
        slide_flow_auto = prs.slides.add_slide(prs.slide_layouts[5])
        slide_flow_auto.shapes.title.text = f"[{date_str}] 自動化後の処理フロー図"
        slide_flow_auto.shapes.title.text_frame.paragraphs[0].font.name = "Meiryo UI"
        
        flow_auto_match = re.search(r'自動化後のフロー図.*?graph TD(.*?)具体的な自動化案とステップ', raw_text, re.DOTALL)
        if flow_auto_match:
            parse_mermaid_to_flowchart(slide_flow_auto, flow_auto_match.group(1))
            
        # 4. 具体的な自動化案とステップ
        steps = extract_steps(soup)
        if steps:
            slide_steps = prs.slides.add_slide(prs.slide_layouts[5])
            slide_steps.shapes.title.text = f"[{date_str}] 具体的な自動化ステップ"
            slide_steps.shapes.title.text_frame.paragraphs[0].font.name = "Meiryo UI"
            
            tb_steps = slide_steps.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
            tf_steps = tb_steps.text_frame
            tf_steps.word_wrap = True
            for step in steps[:4]: # 4つくらいまで
                p = tf_steps.add_paragraph()
                p.text = step['title']
                p.font.name = "Meiryo UI"
                p.font.size = Pt(16)
                p.font.bold = True
                
                p2 = tf_steps.add_paragraph()
                content = step['content'].strip()
                # ログの注釈などは省略
                content = re.sub(r'※現在のログでは.*', '', content, flags=re.DOTALL)
                p2.text = content.strip()[:100] + "..." # 長すぎる場合は省略
                p2.font.name = "Meiryo UI"
                p2.font.size = Pt(14)

    out_path = r"C:\AutoAnalysisLogs\daily_reports\DX_Integrated_Report_Final.pptx"
    prs.save(out_path)
    print("Done:", out_path)

if __name__ == '__main__':
    generate()

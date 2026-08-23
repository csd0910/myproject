import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Title Slide ---
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "長期欠品フラグ更新業務\nDX（Python自動化）時系列詳細比較"
slide.placeholders[1].text = "タイムスタンプと操作ログに基づくAs-Is / To-Be 比較"

# --- Function for chronological comparative slide ---
def add_timeline_slide(title_text, time_str, manual_text, manual_cell, python_text, python_time):
    layout = prs.slide_layouts[5] # Blank with title
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title_text
    
    # Time Box
    time_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.8))
    time_shape.fill.solid()
    time_shape.fill.fore_color.rgb = RGBColor(60, 60, 60)
    time_shape.line.fill.background()
    p = time_shape.text_frame.paragraphs[0]
    p.text = f"【操作時刻】 {time_str}"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    
    # Left Box (As-Is)
    left_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(2.3), Inches(6.0), Inches(4.5))
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = RGBColor(255, 230, 230) 
    left_shape.line.color.rgb = RGBColor(200, 50, 50)
    
    tf = left_shape.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "現在の手作業 (As-Is)"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(150, 0, 0)
    
    p = tf.add_paragraph()
    p.text = manual_text
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.space_before = Pt(15)
    
    p = tf.add_paragraph()
    p.text = f"【セルの変化】\n{manual_cell}"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(100, 50, 50)
    p.space_before = Pt(15)

    # Right Box (To-Be)
    right_shape = slide.shapes.add_shape(1, Inches(6.8), Inches(2.3), Inches(6.0), Inches(4.5))
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = RGBColor(230, 255, 230)
    right_shape.line.color.rgb = RGBColor(50, 200, 50)
    
    tf = right_shape.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "Python自動化後 (To-Be)"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0, 100, 0)
    
    p = tf.add_paragraph()
    p.text = python_text
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.space_before = Pt(15)
    
    p = tf.add_paragraph()
    p.text = f"【処理時間】\n{python_time}"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(50, 100, 50)
    p.font.bold = True
    p.space_before = Pt(15)
    
    # Arrow
    arrow = slide.shapes.add_shape(33, Inches(6.5), Inches(4.0), Inches(0.3), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
    arrow.line.fill.background()

add_timeline_slide(
    "プロセス①：最新データ（CSV）の流し込み",
    "14:48:22 ～ 14:48:29",
    "別ファイルから数万文字の最新データをコピー（Ctrl+C）し、「①長期欠品」シートのA3セル以降に手動で上書き貼り付け（Ctrl+V）。",
    "A3セル等（計2,627箇所）が上書き更新された。（例：840 → 9255 など）",
    "Pythonの pandas.read_csv() で、システム出力した『品目情報_新.csv』などをメモリ上に直接読み込む。人間がファイルを開く必要すらなし。",
    "約 0.2 秒"
)

add_timeline_slide(
    "プロセス②：数十万件のVLOOKUP数式再計算",
    "14:49:14 ～ 14:49:20",
    "Tabキーで右に移動後、Ctrl+Shift+Downでシートの一番下（約60万行）まで選択し、I3～AF3のVLOOKUP数式を一気にオートフィル（Ctrl+V）。\n※マニュアルに「1分ほど固まる」と記載あり。",
    "Z列等の日付やデータが約530万箇所一気に更新。\n※Excelが一時フリーズし、値貼り付けの手間が発生。",
    "Python内で df.merge() を実行し、CSV同士を結合して必要な値（在庫数やランク）を引っ張ってくる処理をバックグラウンドで行う。数式を保存しないため重くならない。",
    "約 2.0 秒"
)

add_timeline_slide(
    "プロセス③：フラグ解除対象の目視検索",
    "14:49:49 ～ 14:50:00",
    "Excel上で Ctrl+F を開き、「J232WF」や「JK4814」といった特定の欠品商品コードを延々と目視検索。AC列～AE列の判定（解除OK等）を見ながら手動で探す。",
    "セルの変化なし（目視と検索による待機・確認時間）",
    "マニュアルの判定ロジック（例：I列=9997なら解除OKなど）をプログラム化。Pythonが全データに対してIF条件を一瞬で通し、該当するリストだけを自動抽出。",
    "約 0.1 秒"
)

add_timeline_slide(
    "プロセス④：手作業による別シートへの移動・微修正",
    "14:50:07 ～ 14:52:50",
    "「解除OK」となった品目をコピーし、「②長期欠品解除」シートへ手動で貼り付け。さらに「社外サイト販売休止」シートの欠品数を矢印キーで移動しながら手修正。",
    "「②長期欠品解除」シートに5,649箇所のデータが追加。\n「社外サイト」のM44セル等が修正（31→56等）。",
    "プロセス③で自動抽出したリストを、「②長期欠品解除」シートに最初から完成した状態で出力（to_excel()）。\n社外サイトの欠品数も差分から自動反映。",
    "約 1.5 秒"
)

out_path = r"C:\Users\フォーレスト026\MyProject\業務自動化分析ツール\AutoAnalysisLogs\daily_reports\Python自動化提案_時系列詳細.pptx"
prs.save(out_path)
print(f"Detailed presentation saved to {out_path}")

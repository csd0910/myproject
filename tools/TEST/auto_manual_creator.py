import time
import cv2
import pyautogui
import numpy as np
from pynput import mouse, keyboard
import uiautomation as auto
from PIL import ImageGrab
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
import os
import datetime
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import ctypes
from ctypes import wintypes

# --- グローバル状態 ---
is_running = False
mouse_listener = None
keyboard_listener = None

last_capture_time = 0
COOLDOWN_SECONDS = 0.5
slides_data = []
TEMP_DIR = "temp_capture_dir"
current_typed_text = ""

settings = {
    "title": "業務マニュアル",
    "author": "作成者名",
    "date": datetime.datetime.now().strftime("%Y/%m/%d"),
    "out_dir": os.path.expanduser("~\\Desktop"),
    "layout": "パターンA (標準:画像左・テキスト右)",
    "font": "メイリオ"
}

if not os.path.exists(TEMP_DIR):
    os.makedirs(TEMP_DIR)


def _record_action(x, y, is_scroll=False):
    global last_capture_time, current_typed_text

    current_time = time.time()
    cooldown = 1.0 if is_scroll else COOLDOWN_SECONDS
    if current_time - last_capture_time < cooldown:
        return
    last_capture_time = current_time

    element_name = "要素"
    top_window = None
    elem_rect = None
    control_type_str = ""
    try:
        with auto.UIAutomationInitializerInThread():
            element = auto.ControlFromCursor()
            if element:
                top_window = element.GetTopLevelControl()
                elem_rect = element.BoundingRectangle
                name = element.Name
                control_type = element.ControlTypeName
                if name:
                    element_name = str(name).strip()
                elif control_type:
                    element_name = str(control_type).replace("Control", "")
                control_type_str = str(control_type)
    except Exception:
        pass

    if len(element_name) > 30:
        element_name = element_name[:28] + "..."

    ignore_words = ["Group", "Window", "Pane", "Document", "要素", "Custom", "TitleBar", "MenuBar"]
    if str(element_name).strip() in ignore_words:
        if is_scroll or "Scroll" in control_type_str or "Thumb" in control_type_str:
            explanation = "画面をスクロール"
        elif "ComboBox" in control_type_str:
            explanation = "プルダウンをクリック"
        else:
            explanation = "画面上をクリック"
    else:
        if is_scroll or "Scroll" in control_type_str or "Thumb" in control_type_str:
            explanation = f"「{element_name}」をスクロール"
        elif "ComboBox" in control_type_str:
            explanation = f"「{element_name}」のプルダウンを選択"
        else:
            explanation = f"「{element_name}」をクリック"

    # アクティブなモニター（マウスがあるモニター）の領域を特定
    def get_monitor_rect(mx, my):
        user32 = ctypes.windll.user32
        h_monitor = user32.MonitorFromPoint(wintypes.POINT(int(mx), int(my)), 1) # 1 = MONITOR_DEFAULTTOPRIMARY
        class MONITORINFO(ctypes.Structure):
            _fields_ = [("cbSize", wintypes.DWORD), ("rcMonitor", wintypes.RECT), ("rcWork", wintypes.RECT), ("dwFlags", wintypes.DWORD)]
        mi = MONITORINFO()
        mi.cbSize = ctypes.sizeof(MONITORINFO)
        if user32.GetMonitorInfoW(h_monitor, ctypes.byref(mi)):
            r = mi.rcMonitor
            return (r.left, r.top, r.right, r.bottom)
        return None

    m_rect = get_monitor_rect(x, y)
    if m_rect:
        left, top, right, bottom = m_rect
        screenshot = ImageGrab.grab(bbox=(left, top, right, bottom), all_screens=True)
        rel_x, rel_y = x - left, y - top
    else:
        # 取得失敗時は全体
        screenshot = ImageGrab.grab(all_screens=True)
        left, top = 0, 0
        rel_x, rel_y = x, y

    frame = np.array(screenshot)
    frame = cv2.cvtColor(frame, cv2.COLOR_RGB2BGR)

    is_same_screen = False
    if len(slides_data) > 0:
        last_frame = slides_data[-1]["raw_image"]
        if last_frame.shape == frame.shape:
            diff = cv2.absdiff(last_frame, frame)
            mean_diff = np.mean(diff)
            if (mean_diff / 255.0) < 0.05:
                is_same_screen = True

    # 実際の要素の大きさに合わせた枠の計算
    if elem_rect and (elem_rect.right - elem_rect.left) > 0 and (elem_rect.bottom - elem_rect.top) > 0:
        pad = 5
        ex = (elem_rect.left - left) - pad
        ey = (elem_rect.top - top) - pad
        ew = (elem_rect.right - elem_rect.left) + pad * 2
        eh = (elem_rect.bottom - elem_rect.top) + pad * 2

        # 大きすぎる場合は制限（画面全体など）
        if ew > frame.shape[1] // 1.5 or eh > frame.shape[0] // 1.5:
            ew, eh = 80, 50
            ex, ey = rel_x - ew // 2, rel_y - eh // 2
    else:
        ew, eh = 80, 50
        ex, ey = rel_x - ew // 2, rel_y - eh // 2

    click_data = {
        "x": rel_x,
        "y": rel_y,
        "rect": (ex, ey, ew, eh),
        "explanation": explanation
    }

    if is_same_screen:
        slides_data[-1]["clicks"].append(click_data)
        print(f"同画面に追加: {explanation}")
    else:
        slides_data.append({
            "raw_image": frame,
            "clicks": [click_data],
            "origin_x": left if top_window else 0,
            "origin_y": top if top_window else 0
        })
        print(f"新画面を取得: {explanation}")

def on_click(x, y, button, pressed):
    global current_typed_text
    if not is_running:
        return

    if pressed and button == mouse.Button.left:
        # 次のクリックが行われる前に文字があれば手前の操作に追記
        if current_typed_text.strip() and len(slides_data) > 0:
            last_click = slides_data[-1]["clicks"][-1]
            if "（入力:" not in last_click["explanation"]:
                last_click["explanation"] += f" （入力: {current_typed_text.strip()}）"
        current_typed_text = ""

        _record_action(x, y, is_scroll=False)

def on_scroll(x, y, dx, dy):
    global current_typed_text
    if not is_running:
        return

    # スクロール時も文字があれば手前の操作に追記
    if current_typed_text.strip() and len(slides_data) > 0:
        last_click = slides_data[-1]["clicks"][-1]
        if "（入力:" not in last_click["explanation"]:
            last_click["explanation"] += f" （入力: {current_typed_text.strip()}）"
    current_typed_text = ""

    _record_action(x, y, is_scroll=True)

def show_edit_window():
    import PIL.ImageTk
    from PIL import Image, ImageDraw

    edit_win = tk.Toplevel(app)
    edit_win.title("手順の確認・削除")
    edit_win.geometry("1200x850")

    # [X]ボタンで閉じられた時にメインUIに戻す
    edit_win.protocol("WM_DELETE_WINDOW", lambda: [edit_win.destroy(), app.deiconify()])

    tk.Label(edit_win, text="不要な操作を選択して「削除」できます。右側にプレビューが表示されます。\n削除・修正すると番号は自動で振り直されます。", font=("Meiryo", 10)).pack(pady=5)

    flat_clicks = []
    for s_idx, slide_info in enumerate(slides_data):
        for click in slide_info["clicks"]:
            flat_clicks.append((s_idx, click, click["explanation"]))

    main_frame = tk.Frame(edit_win)
    main_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)

    left_frame = tk.Frame(main_frame)
    left_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

    right_frame = tk.Frame(main_frame, width=820)
    right_frame.pack(side=tk.RIGHT, fill=tk.BOTH, expand=True, padx=10)

    listbox = tk.Listbox(left_frame, selectmode=tk.EXTENDED, font=("Meiryo", 11))
    listbox.pack(fill=tk.BOTH, expand=True)

    preview_label = tk.Label(right_frame, text="画像プレビュー\n(左のリストを選択してください)", bg="#e0e0e0", font=("Meiryo", 10))
    preview_label.pack(fill=tk.BOTH, expand=True)
    preview_label.image_ref = None

    edit_frame = tk.Frame(right_frame, relief=tk.RIDGE, bd=2, padx=5, pady=5)
    edit_frame.pack(fill=tk.X, side=tk.BOTTOM, pady=(5, 0))

    tk.Label(edit_frame, text="出力文言を編集:", font=("Meiryo", 10)).pack(side=tk.TOP, anchor="w")

    desc_var = tk.StringVar(value="(未選択)")
    desc_entry = tk.Entry(edit_frame, textvariable=desc_var, font=("Meiryo", 11, "bold"))
    desc_entry.pack(fill=tk.X, side=tk.TOP, pady=2)

    current_idx_ref = {"idx": -1}

    def on_update_desc():
        idx = current_idx_ref["idx"]
        if idx >= 0 and idx < len(flat_clicks):
            new_val = desc_var.get()
            s_idx, click, expl = flat_clicks[idx]
            click["explanation"] = new_val # slides_data側の参照を直接更新
            flat_clicks[idx] = (s_idx, click, new_val)

            # リストの表記を更新
            listbox.delete(idx)
            listbox.insert(idx, f"{idx+1}. {new_val}")
            listbox.selection_set(idx)
            messagebox.showinfo("更 新", "文言を修正しました！\n（PowerPointにも反映されます）", parent=edit_win)

    tk.Button(edit_frame, text="変更を保存", font=("Meiryo", 10), bg="#e1f0fa", command=on_update_desc).pack(side=tk.RIGHT, pady=2)

    def on_select(evt):
        sel = listbox.curselection()
        if not sel:
            return
        idx = sel[0]
        if idx >= len(flat_clicks): return
        s_idx, click, expl = flat_clicks[idx]
        frame = slides_data[s_idx]["raw_image"]
        # Convert BGR to RGB
        rgb_frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
        pil_img = Image.fromarray(rgb_frame)

        draw = ImageDraw.Draw(pil_img)

        cx, cy = int(click["x"]), int(click["y"])
        if "rect" in click:
            ex, ey, ew, eh = click["rect"]
            draw.rectangle([ex, ey, ex+ew, ey+eh], outline="red", width=6)
            # 左上の角を目印の基準点にする
            px, py = ex, ey
        else:
            rw, rh = 20, 15
            draw.rectangle([cx-rw, cy-rh, cx+rw, cy+rh], outline="red", width=6)
            # 左上の角を目印の基準点にする
            px, py = cx-rw, cy-rh

        # わかりやすいように半透明の赤い円を描く（枠の左上固定）
        draw.ellipse([px-20, py-20, px+20, py+20], fill="red")

        # Resize image for preview
        img_w, img_h = pil_img.size
        max_w, max_h = 800, 600
        scale = min(max_w / max(1, img_w), max_h / max(1, img_h))
        new_w, new_h = max(1, int(img_w * scale)), max(1, int(img_h * scale))
        pil_img = pil_img.resize((new_w, new_h), Image.LANCZOS)

        photo = PIL.ImageTk.PhotoImage(pil_img)
        preview_label.config(image=photo, text="")
        preview_label.image_ref = photo

        # 選択されたインデックスと文言をテキストボックスに反映
        current_idx_ref["idx"] = idx
        desc_var.set(expl)

    listbox.bind('<<ListboxSelect>>', on_select)

    def refresh_list():
        listbox.delete(0, tk.END)
        for i, (_, _, expl) in enumerate(flat_clicks):
            listbox.insert(tk.END, f"{i+1}. {expl}")
        preview_label.config(image='', text="画像プレビュー\n(左のリストを選択してください)")
        desc_var.set("(未選択)")
        current_idx_ref["idx"] = -1

    refresh_list()

    def on_delete():
        selected = list(listbox.curselection())
        selected.reverse()
        for i in selected:
            flat_clicks.pop(i)
        refresh_list()

    tk.Button(left_frame, text="選択した工程を削除", font=("Meiryo", 10), command=on_delete).pack(pady=5)

    def on_generate():
        new_slides = []
        for s_idx, slide_info in enumerate(slides_data):
            new_clicks = [c for s, c, e in flat_clicks if s == s_idx]
            if new_clicks:
                new_slides.append({
                    "raw_image": slide_info["raw_image"],
                    "clicks": new_clicks,
                    "origin_x": slide_info["origin_x"],
                    "origin_y": slide_info["origin_y"]
                })
        generate_pptx(new_slides, parent_win=edit_win)

    tk.Button(left_frame, text="この内容でPowerPointを出力", bg="#0078D7", fg="white", font=("Meiryo", 12, "bold"), command=on_generate).pack(pady=10)

    def on_close_to_main():
        edit_win.destroy()
        app.deiconify()

    tk.Button(left_frame, text="設定画面に戻る (記録をリセット)", font=("Meiryo", 10), command=on_close_to_main).pack(pady=5)

def on_press(key):
    global is_running, current_typed_text
    if key == keyboard.Key.f12 and is_running:
        # 終了時も、打った文字があれば追記
        if current_typed_text.strip() and len(slides_data) > 0:
            last_click = slides_data[-1]["clicks"][-1]
            if "（入力:" not in last_click["explanation"]:
                last_click["explanation"] += f" （入力: {current_typed_text.strip()}）"
        current_typed_text = ""

        print("F12キー検知。プレビュー・編集画面を開きます...")
        is_running = False
        app.after(0, show_edit_window)
        return False

    if is_running:
        try:
            if hasattr(key, 'char') and key.char is not None:
                current_typed_text += key.char
            elif key == keyboard.Key.space:
                current_typed_text += " "
            elif key == keyboard.Key.backspace:
                current_typed_text = current_typed_text[:-1]
        except Exception:
            pass

def start_capture():
    global is_running, mouse_listener, keyboard_listener, slides_data
    slides_data.clear()
    is_running = True
    print("キャプチャ開始。F12で終了。")
    mouse_listener = mouse.Listener(on_click=on_click, on_scroll=on_scroll)
    keyboard_listener = keyboard.Listener(on_press=on_press)
    mouse_listener.start()
    keyboard_listener.start()

def generate_pptx(data, parent_win=None):
    try:
        if not data:
            if parent_win:
                messagebox.showinfo("情報", "出力する操作がありません。", parent=parent_win)
            else:
                root = tk.Tk()
                root.withdraw()
                messagebox.showinfo("情報", "キャプチャされた画像がありません。")
                root.destroy()
                app.deiconify() # UIを再表示
            return

        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        title_slide_layout = prs.slide_layouts[0]
        blank_slide_layout = prs.slide_layouts[6]
        font_name = settings["font"]

        def apply_font(run, size_pt, bold=False, color_rgb=RGBColor(0,0,0)):
            run.font.name = font_name
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            run.font.color.rgb = color_rgb

        def add_white_bg(slide):
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)

        # ---------------- 1. 表紙 ----------------
        slide_title = prs.slides.add_slide(title_slide_layout)
        title_shape = slide_title.shapes.title
        subtitle_shape = slide_title.placeholders[1]

        title_shape.text = settings["title"]
        apply_font(title_shape.text_frame.paragraphs[0].runs[0], 54, bold=True)

        subtitle_text = f"作成者: {settings['author']}\n作成日: {settings['date']}"
        subtitle_shape.text = subtitle_text
        for p in subtitle_shape.text_frame.paragraphs:
            for r in p.runs:
                apply_font(r, 24)

        # 全操作テキストの抽出
        all_operations = []
        for slide_info in data:
            for click in slide_info["clicks"]:
                all_operations.append(click["explanation"])

        # ---------------- 2. 目次 (Index) ----------------
        slide_index = prs.slides.add_slide(blank_slide_layout)
        add_white_bg(slide_index)
        tx_box = slide_index.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1))
        tx_box.text_frame.text = "目次 (Index)"
        apply_font(tx_box.text_frame.paragraphs[0].runs[0], 40, bold=True)

        content_box = slide_index.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(6))
        tf = content_box.text_frame
        tf.word_wrap = True

        limit_idx = min(len(all_operations), 15)
        for i in range(limit_idx):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = f"{i+1}. 「{all_operations[i]}」の操作"
            apply_font(p.runs[0], 20)

        if len(all_operations) > 15:
            p = tf.add_paragraph()
            p.text = "他..."
            apply_font(p.runs[0], 20)

        # ---------------- 3. 業務フロー・概要 ----------------
        slide_flow = prs.slides.add_slide(blank_slide_layout)
        add_white_bg(slide_flow)
        tx_box = slide_flow.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1))
        tx_box.text_frame.text = "業務フロー・概要"
        apply_font(tx_box.text_frame.paragraphs[0].runs[0], 40, bold=True)

        # 矢印でフロー表現（図形でグラフィカルに配置）
        start_x = Inches(1)
        start_y = Inches(2.2)

        # アイテム数によってサイズと改行数を動的に計算する
        total_items = len(all_operations)
        if total_items <= 25:
            box_width = Inches(2.6)
            box_height = Inches(0.8)
            font_size = 14
        elif total_items <= 50:
            box_width = Inches(1.8)
            box_height = Inches(0.6)
            font_size = 10
        else:
            box_width = Inches(1.3)
            box_height = Inches(0.4)
            font_size = 8

        x_gap = Inches(0.2)
        y_gap = Inches(0.2) if total_items > 25 else Inches(0.5)

        cur_x = start_x
        cur_y = start_y

        for j, op in enumerate(all_operations):
            shape = slide_flow.shapes.add_shape(
                MSO_SHAPE.CHEVRON, cur_x, cur_y, box_width, box_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0, 120, 215) # 青色
            shape.line.color.rgb = RGBColor(0, 120, 215)

            tf_shape = shape.text_frame
            tf_shape.word_wrap = True
            tf_shape.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p_shape = tf_shape.paragraphs[0]

            p_shape.text = f"{j+1}. {op}"
            p_shape.alignment = PP_ALIGN.CENTER

            # PowerShellの自動縮小だけでなく、自前でも長すぎる場合は少し文字を小さくしておく
            dynamic_font_size = font_size
            if len(p_shape.text) > 12:
                dynamic_font_size = max(6, int(font_size * (12 / len(p_shape.text))))

            apply_font(p_shape.runs[0], dynamic_font_size, bold=True, color_rgb=RGBColor(255, 255, 255))

            cur_x += box_width + x_gap
            if cur_x + box_width > Inches(15): # 右端にはみ出しそうなら改行
                cur_x = start_x
                cur_y += box_height + y_gap

        # ---------------- 4. 個別操作スライド ----------------
        layout_type = settings["layout"]

        for i, slide_info in enumerate(data):
            slide = prs.slides.add_slide(blank_slide_layout)
            add_white_bg(slide)

            frame = slide_info["raw_image"]
            ih, iw, _ = frame.shape
            explanations_text = []

            for idx, click in enumerate(slide_info["clicks"]):
                circled_numbers = ["①","②","③","④","⑤","⑥","⑦","⑧","⑨","⑩","⑪","⑫","⑬","⑭","⑮"]
                mark = circled_numbers[idx] if idx < len(circled_numbers) else f"({idx+1})"
                explanations_text.append(f"{mark} {click['explanation']}")

            # 加工のないそのままの綺麗な画像を保存
            img_filename = os.path.join(TEMP_DIR, f"slide_out_{i}.jpg")
            cv2.imwrite(img_filename, frame, [int(cv2.IMWRITE_JPEG_QUALITY), 80])

            # レイアウトに基づく最大サイズの決定
            margin = Inches(0.2)
            if "画像左" in layout_type:
                max_w = Inches(9.5)
                max_h = Inches(7.5)
                img_left = margin
                img_top = Inches(1)
                tx_left_base = img_left + max_w + Inches(0.3)
                tx_top = img_top
                tx_width = Inches(5.5)
                tx_height = Inches(7.5)
            elif "画像上" in layout_type:
                max_w = Inches(10)
                max_h = Inches(5.5)
                img_left = Inches(3) # 余白
                img_top = Inches(1)
                tx_left_base = margin
                tx_top = img_top + max_h + Inches(0.2)
                tx_width = Inches(15.6)
                tx_height = Inches(2)
            elif "画像のみ" in layout_type:
                max_w = Inches(15.6)
                max_h = Inches(8.6)
                img_left = Inches(0.2)
                img_top = Inches(0.2)
                tx_left_base, tx_top, tx_width, tx_height = 0, 0, 0, 0
            else: # パターンC: 画像大きく＋下テキスト
                max_w = Inches(14)
                max_h = Inches(6.0)
                img_left = Inches(1)
                img_top = Inches(0.2)
                tx_left_base = Inches(1)
                tx_top = img_top + max_h + Inches(0.2)
                tx_width = Inches(14)
                tx_height = Inches(2)

            # 画像のスケーリング（大きい場合はアスペクト比を維持して枠内に収める）
            scale_w = max_w / iw
            scale_h = max_h / ih
            scale = min(scale_w, scale_h)
            render_w = int(iw * scale)
            render_h = int(ih * scale)

            if "画像左" in layout_type:
                tx_left = img_left + render_w + Inches(0.3)
            else:
                tx_left = tx_left_base

            # 見出し・画像の配置
            if "画像のみ" not in layout_type:
                tx_title = slide.shapes.add_textbox(margin, margin, Inches(15), Inches(0.8))
                p_title = tx_title.text_frame.paragraphs[0]
                p_title.text = f"【ここに大見出しを入力 - 手順{i+1}】"
                apply_font(p_title.runs[0], 32, bold=True)

            pic = slide.shapes.add_picture(img_filename, img_left, img_top, width=render_w, height=render_h)

            # 画像の上に直接、PowerPointの編集可能な図形として赤枠を描画
            render_l = pic.left
            render_t = pic.top

            # 座標が近いクリックは同じ赤枠にまとめる処理
            click_groups = []
            for idx, click in enumerate(slide_info["clicks"]):
                cx = click["x"]
                cy = click["y"]
                grouped = False
                for grp in click_groups:
                    if abs(cx - grp["x"]) < 30 and abs(cy - grp["y"]) < 30: # 30px以内は同一グループ
                        grp["indices"].append(idx)
                        grouped = True
                        break
                if not grouped:
                    click_groups.append({"x": cx, "y": cy, "indices": [idx]})

            for grp in click_groups:
                cx = grp["x"]
                cy = grp["y"]
                first_idx = grp["indices"][0]
                first_click = slide_info["clicks"][first_idx]

                if "rect" in first_click:
                    ex, ey, ew, eh = first_click["rect"]
                    rect_l = render_l + ex * scale
                    rect_t = render_t + ey * scale
                    rect_w = ew * scale
                    rect_h = eh * scale
                else:
                    pt_x = render_l + cx * scale
                    pt_y = render_t + cy * scale
                    rect_w = Inches(0.8)
                    rect_h = Inches(0.4)
                    rect_l = pt_x - rect_w / 2
                    rect_t = pt_y - rect_h / 2

                # PowerPointの矩形図形（赤枠・塗りつぶしなし）
                shape_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rect_l, rect_t, rect_w, rect_h)
                shape_rect.fill.background()
                shape_rect.line.color.rgb = RGBColor(255, 0, 0)
                shape_rect.line.width = Pt(3)

                # 重複時の文字連結（①⑧など）
                circled_numbers = ["①","②","③","④","⑤","⑥","⑦","⑧","⑨","⑩","⑪","⑫","⑬","⑭","⑮"]
                text_marks = "".join([circled_numbers[idx] if idx < len(circled_numbers) else f"({idx+1})" for idx in grp["indices"]])

                # PowerPointのラベル図形（赤色・白文字）
                char_count = len(text_marks)
                oval_w = Inches(0.35 + (char_count - 1) * 0.25)
                oval_h = Inches(0.35)
                shape_oval = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE if char_count > 1 else MSO_SHAPE.OVAL,
                    rect_l, rect_t - oval_h, oval_w, oval_h
                )
                shape_oval.fill.solid()
                shape_oval.fill.fore_color.rgb = RGBColor(255, 0, 0)
                shape_oval.line.color.rgb = RGBColor(255, 0, 0)

                tf_oval = shape_oval.text_frame
                tf_oval.word_wrap = False
                tf_oval.margin_left, tf_oval.margin_right = 0, 0
                tf_oval.margin_top, tf_oval.margin_bottom = 0, 0
                p_oval = tf_oval.paragraphs[0]
                p_oval.text = text_marks
                p_oval.alignment = PP_ALIGN.CENTER
                apply_font(p_oval.runs[0], 14, bold=True, color_rgb=RGBColor(255, 255, 255))

            if "画像のみ" not in layout_type:
                textbox = slide.shapes.add_textbox(tx_left, tx_top, tx_width, tx_height)
                tf = textbox.text_frame
                tf.word_wrap = True

                p = tf.paragraphs[0]
                p.text = "【具体的な説明をここに入力してください】\n\n" + "\n".join(explanations_text)
                for r in p.runs:
                    apply_font(r, 22)

        out_name = "AutoManual_Output_" + datetime.datetime.now().strftime("%Y%m%d_%H%M%S") + ".pptx"
        output_path = os.path.join(settings["out_dir"], out_name)
        prs.save(output_path)

        try:
            for f in os.listdir(TEMP_DIR):
                os.remove(os.path.join(TEMP_DIR, f))
        except Exception:
            pass

        print(f"作成完了: {output_path}")

        # UIを再開・通知
        if parent_win is not None:
            messagebox.showinfo("完了", f"PowerPointの生成が完了しました！\n引き続き編集するか、設定画面に戻れます。\n\n保存先:\n{output_path}", parent=parent_win)
        else:
            root = tk.Tk()
            root.withdraw()
            messagebox.showinfo("完了", f"生成が完了しました！\n\n保存先:\n{output_path}")
            root.destroy()
            app.deiconify()

    except Exception as e:
        print(f"PPTX生成エラー: {e}")
        import traceback
        traceback.print_exc()

# --- GUIクラス ---
class AppUI(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("自動マニュアル生成ツール 設定")
        self.geometry("450x550")
        self.configure(padx=20, pady=20)
        self.create_widgets()

    def create_widgets(self):
        label_font = ("Meiryo", 10)
        title_font = ("Meiryo", 14, "bold")

        # タイトル
        tk.Label(self, text="自動マニュアル作成", font=title_font).pack(pady=(0,20))

        # Title
        frame_title = tk.Frame(self)
        frame_title.pack(fill=tk.X, pady=5)
        tk.Label(frame_title, text="表紙タイトル:", width=20, anchor="w", font=label_font).pack(side=tk.LEFT)
        self.var_title = tk.StringVar(value=settings["title"])
        tk.Entry(frame_title, textvariable=self.var_title, width=28, font=label_font).pack(side=tk.LEFT)

        # Author
        frame_author = tk.Frame(self)
        frame_author.pack(fill=tk.X, pady=5)
        tk.Label(frame_author, text="作成者名:", width=20, anchor="w", font=label_font).pack(side=tk.LEFT)
        self.var_author = tk.StringVar(value=settings["author"])
        tk.Entry(frame_author, textvariable=self.var_author, width=28, font=label_font).pack(side=tk.LEFT)

        # Date
        frame_date = tk.Frame(self)
        frame_date.pack(fill=tk.X, pady=5)
        tk.Label(frame_date, text="作成日:", width=20, anchor="w", font=label_font).pack(side=tk.LEFT)
        self.var_date = tk.StringVar(value=settings["date"])
        tk.Entry(frame_date, textvariable=self.var_date, width=28, font=label_font).pack(side=tk.LEFT)

        # Output Dir
        frame_out = tk.Frame(self)
        frame_out.pack(fill=tk.X, pady=5)
        tk.Label(frame_out, text="保存先フォルダ:", width=20, anchor="w", font=label_font).pack(side=tk.LEFT)
        self.var_out = tk.StringVar(value=settings["out_dir"])
        tk.Entry(frame_out, textvariable=self.var_out, width=18, font=label_font).pack(side=tk.LEFT)
        tk.Button(frame_out, text="参照", command=self.select_dir, font=label_font).pack(side=tk.LEFT, padx=5)

        # Layout
        frame_layout = tk.Frame(self)
        frame_layout.pack(fill=tk.X, pady=5)
        tk.Label(frame_layout, text="スライドレイアウト:", width=20, anchor="w", font=label_font).pack(side=tk.LEFT)
        self.var_layout = tk.StringVar(value=settings["layout"])
        layouts = ["パターンA (標準:画像左・テキスト右)", "パターンB (画像上・テキスト下)", "パターンC (大きな画像＋下テキスト)", "パターンD (画像のみ)"]
        cb_layout = ttk.Combobox(frame_layout, textvariable=self.var_layout, values=layouts, state="readonly", width=26, font=label_font)
        cb_layout.pack(side=tk.LEFT)

        # Font
        frame_font = tk.Frame(self)
        frame_font.pack(fill=tk.X, pady=5)
        tk.Label(frame_font, text="PowerPointフォント:", width=20, anchor="w", font=label_font).pack(side=tk.LEFT)
        self.var_font = tk.StringVar(value=settings["font"])
        fonts = ["メイリオ", "游ゴシック", "BIZ UDPゴシック", "Yu Gothic UI", "MS Pゴシック"]
        cb_font = ttk.Combobox(frame_font, textvariable=self.var_font, values=fonts, state="readonly", width=26, font=label_font)
        cb_font.pack(side=tk.LEFT)

        # 説明文
        info_text = (
            "【使い方】\n"
            "1. 上記を設定し「キャプチャ開始」をクリック。\n"
            "2. この画面が消え、記録モードに入ります。\n"
            "3. 説明したい場所を左クリックしていくと\n"
            "   画面と操作内容が自動保存されます。\n"
            "4. 終了時はキーボードの【F12】を押してください。\n"
            "   自動でPowerPointが生成されます。"
        )
        tk.Label(self, text=info_text, fg="#555", justify=tk.LEFT, font=("Meiryo", 9)).pack(pady=15)

        # Start Button
        self.btn_start = tk.Button(self, text="キャプチャ開始 (記録モードへ)", font=("Meiryo", 12, "bold"), bg="#0078D7", fg="white", command=self.on_start)
        self.btn_start.pack(fill=tk.X, pady=10, ipady=5)

    def select_dir(self):
        d = filedialog.askdirectory(initialdir=self.var_out.get())
        if d:
            self.var_out.set(d)

    def on_start(self):
        settings["title"] = self.var_title.get()
        settings["author"] = self.var_author.get()
        settings["date"] = self.var_date.get()
        settings["out_dir"] = self.var_out.get()
        settings["layout"] = self.var_layout.get()
        settings["font"] = self.var_font.get()

        self.withdraw() # ウィンドウを隠す
        start_capture()

if __name__ == "__main__":
    app = AppUI()
    app.mainloop()

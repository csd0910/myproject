import os
import json
import threading
import tkinter as tk
from tkinter import ttk, filedialog, messagebox

try:
    import openpyxl
except ImportError:
    openpyxl = None
try:
    import docx
except ImportError:
    docx = None
try:
    import pptx
except ImportError:
    pptx = None
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

# 設定ファイル名と制限
SETTINGS_FILE = "search_settings.json"
MAX_FILE_SIZE_MB = 10
MAX_SIZE_BYTES = MAX_FILE_SIZE_MB * 1024 * 1024

class FileSearchApp:
    def __init__(self, root):
        self.root = root
        self.root.title("ファイル内テキスト検索ツール")
        self.root.geometry("800x650")
        
        self.settings = self.load_settings()
        self.is_searching = False
        self.cancel_flag = False
        
        self.create_widgets()
        
    def load_settings(self):
        if os.path.exists(SETTINGS_FILE):
            try:
                with open(SETTINGS_FILE, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except Exception:
                pass
        return {"last_dir": "", "history": []}
        
    def save_settings(self):
        try:
            with open(SETTINGS_FILE, 'w', encoding='utf-8') as f:
                json.dump(self.settings, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"設定保存エラー: {e}")

    def create_widgets(self):
        top_frame = ttk.Frame(self.root, padding=10)
        top_frame.pack(fill=tk.X)
        
        ttk.Label(top_frame, text="対象フォルダ:").grid(row=0, column=0, sticky=tk.W, pady=5)
        self.dir_var = tk.StringVar(value=self.settings.get("last_dir", ""))
        self.dir_entry = ttk.Entry(top_frame, textvariable=self.dir_var, width=60)
        self.dir_entry.grid(row=0, column=1, padx=5, pady=5)
        ttk.Button(top_frame, text="参照", command=self.browse_dir).grid(row=0, column=2, pady=5)
        
        ttk.Label(top_frame, text="検索キーワード:").grid(row=1, column=0, sticky=tk.W, pady=5)
        self.keyword_var = tk.StringVar()
        self.keyword_combo = ttk.Combobox(top_frame, textvariable=self.keyword_var, width=65, font=("Meiryo UI", 11), values=self.settings.get("history", []))
        self.keyword_combo.grid(row=1, column=1, padx=5, pady=(5, 0), sticky=tk.W)
        
        ttk.Label(top_frame, text="※複数検索する場合はカンマ（,）区切りで入力してください", foreground="gray", font=("Meiryo UI", 8)).grid(row=2, column=1, sticky=tk.W, padx=5, pady=(0, 10))
        
        # AND / OR 選択
        self.search_cond_var = tk.StringVar(value="AND")
        radio_frame = ttk.Frame(top_frame)
        radio_frame.grid(row=3, column=1, pady=5, sticky=tk.W)
        ttk.Radiobutton(radio_frame, text="すべて含む (AND)", variable=self.search_cond_var, value="AND").pack(side=tk.LEFT, padx=(0, 15))
        ttk.Radiobutton(radio_frame, text="どれかを含む (OR)", variable=self.search_cond_var, value="OR").pack(side=tk.LEFT)
        
        btn_frame = ttk.Frame(top_frame)
        btn_frame.grid(row=4, column=1, pady=10, sticky=tk.W)
        self.search_btn = ttk.Button(btn_frame, text="検索開始", command=self.start_search)
        self.search_btn.pack(side=tk.LEFT, padx=5)
        self.cancel_btn = ttk.Button(btn_frame, text="キャンセル", command=self.cancel_search, state=tk.DISABLED)
        self.cancel_btn.pack(side=tk.LEFT, padx=5)
        
        self.status_var = tk.StringVar(value="待機中...")
        ttk.Label(top_frame, textvariable=self.status_var).grid(row=5, column=1, sticky=tk.W)
        
        self.result_text = tk.Text(self.root, wrap=tk.NONE, font=("Meiryo UI", 9))
        self.result_text.pack(fill=tk.BOTH, expand=True, padx=10, pady=(0, 10))
        
        # ハイライト用・リンク用のタグ設定
        self.result_text.tag_config("highlight", background="yellow", foreground="black")
        self.result_text.tag_config("link", foreground="blue", underline=1)
        
        # バッジ用のタグ設定
        self.result_text.tag_config("badge_xlsx", background="#1D6F42", foreground="white")
        self.result_text.tag_config("badge_docx", background="#2B579A", foreground="white")
        self.result_text.tag_config("badge_pptx", background="#D24726", foreground="white")
        self.result_text.tag_config("badge_pdf", background="#F40F02", foreground="white")
        self.result_text.tag_config("badge_csv", background="#5C5C5C", foreground="white")
        self.result_text.tag_config("badge_txt", background="#808080", foreground="white")
        
        # リンクのイベント設定
        self.result_text.bind("<Button-1>", self.on_text_click)
        self.result_text.tag_bind("link", "<Enter>", lambda e: self.result_text.config(cursor="hand2"))
        self.result_text.tag_bind("link", "<Leave>", lambda e: self.result_text.config(cursor=""))
        
        ysb = ttk.Scrollbar(self.result_text, orient=tk.VERTICAL, command=self.result_text.yview)
        ysb.pack(side=tk.RIGHT, fill=tk.Y)
        xsb = ttk.Scrollbar(self.result_text, orient=tk.HORIZONTAL, command=self.result_text.xview)
        xsb.pack(side=tk.BOTTOM, fill=tk.X)
        self.result_text.configure(yscrollcommand=ysb.set, xscrollcommand=xsb.set)

    def on_text_click(self, event):
        index = self.result_text.index(f"@{event.x},{event.y}")
        tags = self.result_text.tag_names(index)
        for tag in tags:
            if tag.startswith("file:"):
                filepath = tag.split("file:", 1)[1]
                try:
                    os.startfile(filepath)
                except Exception as e:
                    messagebox.showerror("エラー", f"ファイルを開けませんでした:\n{e}")
                break

    def browse_dir(self):
        d = filedialog.askdirectory(initialdir=self.dir_var.get())
        if d:
            self.dir_var.set(d)
            
    def start_search(self):
        target_dir = self.dir_var.get().strip()
        raw_keyword = self.keyword_var.get().strip()
        
        if not target_dir or not os.path.isdir(target_dir):
            messagebox.showwarning("警告", "有効なフォルダを選択してください。")
            return
        if not raw_keyword:
            messagebox.showwarning("警告", "検索キーワードを入力してください。")
            return
            
        # カンマ区切り処理（全角・半角対応）
        raw_keyword_fmt = raw_keyword.replace('、', ',')
        keywords = [k.strip() for k in raw_keyword_fmt.split(',') if k.strip()]
        
        if not keywords:
            messagebox.showwarning("警告", "有効なキーワードが見つかりません。")
            return
            
        self.settings["last_dir"] = target_dir
        history = self.settings.get("history", [])
        if raw_keyword in history:
            history.remove(raw_keyword)
        history.insert(0, raw_keyword)
        self.settings["history"] = history[:10]
        self.keyword_combo['values'] = self.settings["history"]
        self.save_settings()
        
        self.result_text.delete(1.0, tk.END)
        self.is_searching = True
        self.cancel_flag = False
        self.search_btn.config(state=tk.DISABLED)
        self.cancel_btn.config(state=tk.NORMAL)
        
        cond = self.search_cond_var.get()
        threading.Thread(target=self.search_worker, args=(target_dir, keywords, cond), daemon=True).start()

    def cancel_search(self):
        self.cancel_flag = True
        self.status_var.set("キャンセルしています...")

    def log_result(self, msg):
        self.result_text.insert(tk.END, msg + "\n")
        self.result_text.see(tk.END)
        
    def get_badge_info(self, ext):
        if ext == '.xlsx': return " Excel ", "badge_xlsx"
        elif ext == '.docx': return " Word ", "badge_docx"
        elif ext == '.pptx': return " PowerPoint ", "badge_pptx"
        elif ext == '.pdf': return " PDF ", "badge_pdf"
        elif ext == '.csv': return " CSV ", "badge_csv"
        else: return " Text ", "badge_txt"

    def log_result_formatted(self, filepath, snippet, keywords, ext):
        file_tag = f"file:{filepath}"
        badge_text, badge_tag = self.get_badge_info(ext)
        
        self.result_text.insert(tk.END, "【発見】 ")
        self.result_text.insert(tk.END, badge_text, badge_tag)
        self.result_text.insert(tk.END, " ")
        self.result_text.insert(tk.END, f"{filepath}\n", ("link", file_tag))
        
        self.result_text.insert(tk.END, "  内容: ")
        
        # Snippetを挿入し、複数キーワードすべてをハイライト処理
        start_idx = self.result_text.index(tk.INSERT)
        self.result_text.insert(tk.END, snippet)
        
        for k in keywords:
            pos = 0
            while True:
                idx = snippet.find(k, pos)
                if idx == -1:
                    break
                highlight_start = f"{start_idx} + {idx} chars"
                highlight_end = f"{highlight_start} + {len(k)} chars"
                self.result_text.tag_add("highlight", highlight_start, highlight_end)
                pos = idx + len(k)
                
        self.result_text.insert(tk.END, "\n" + "-" * 50 + "\n")
        self.result_text.see(tk.END)

    def search_worker(self, target_dir, keywords, cond):
        try:
            self.status_var.set("検索中...")
            count = 0
            found_count = 0
            
            for root_dir, dirs, files in os.walk(target_dir):
                if self.cancel_flag:
                    break
                for file in files:
                    if self.cancel_flag:
                        break
                        
                    ext = os.path.splitext(file)[1].lower()
                    if ext not in ['.txt', '.csv', '.xlsx', '.docx', '.pptx', '.pdf']:
                        continue
                        
                    filepath = os.path.join(root_dir, file)
                    count += 1
                    
                    try:
                        size = os.path.getsize(filepath)
                        if size > MAX_SIZE_BYTES:
                            continue 
                            
                        snippet = self.search_in_file(filepath, ext, keywords, cond)
                        if snippet:
                            found_count += 1
                            self.root.after(0, self.log_result_formatted, filepath, snippet.strip(), keywords, ext)
                            
                    except Exception:
                        pass

            if self.cancel_flag:
                self.root.after(0, self.status_var.set, f"キャンセルされました。 (検索したファイル数: {count}, 該当: {found_count})")
            else:
                self.root.after(0, self.status_var.set, f"検索完了！ (検索したファイル数: {count}, 該当: {found_count})")
                if found_count == 0:
                    self.root.after(0, self.log_result, "キーワードは見つかりませんでした。")
                    
        finally:
            self.root.after(0, self.search_btn.config, {"state": tk.NORMAL})
            self.root.after(0, self.cancel_btn.config, {"state": tk.DISABLED})
            self.is_searching = False

    def _check_and_get_snippet(self, text, keywords, cond):
        if not text:
            return None
            
        if cond == "AND":
            for k in keywords:
                if k not in text:
                    return None
            # 全て存在する場合、最初のキーワード付近をスニペットとする
            return self._extract_snippet(text, keywords[0])
        else:
            # OR検索
            for k in keywords:
                if k in text:
                    return self._extract_snippet(text, k)
            return None
            
    def _extract_snippet(self, text, matched_keyword):
        idx = text.find(matched_keyword)
        if idx != -1:
            start = max(0, idx - 30)
            end = min(len(text), idx + len(matched_keyword) + 30)
            return text[start:end].replace('\n', ' ')
        return None

    def search_in_file(self, filepath, ext, keywords, cond):
        if ext in ['.txt', '.csv']:
            return self.search_text(filepath, keywords, cond)
        elif ext == '.xlsx' and openpyxl:
            return self.search_xlsx(filepath, keywords, cond)
        elif ext == '.docx' and docx:
            return self.search_docx(filepath, keywords, cond)
        elif ext == '.pptx' and pptx:
            return self.search_pptx(filepath, keywords, cond)
        elif ext == '.pdf' and PyPDF2:
            return self.search_pdf(filepath, keywords, cond)
        return None

    def search_text(self, filepath, keywords, cond):
        encodings = ['utf-8', 'cp932', 'shift_jis']
        for enc in encodings:
            try:
                with open(filepath, 'r', encoding=enc, errors='replace') as f:
                    content = f.read()
                    return self._check_and_get_snippet(content, keywords, cond)
            except Exception:
                continue
        return None

    def search_xlsx(self, filepath, keywords, cond):
        try:
            wb = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
            parts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell and isinstance(cell, str):
                            parts.append(cell)
            full_text = " ".join(parts)
            return self._check_and_get_snippet(full_text, keywords, cond)
        except Exception:
            return None

    def search_docx(self, filepath, keywords, cond):
        try:
            doc = docx.Document(filepath)
            full_text = " ".join([para.text for para in doc.paragraphs])
            return self._check_and_get_snippet(full_text, keywords, cond)
        except Exception:
            return None

    def search_pptx(self, filepath, keywords, cond):
        try:
            prs = pptx.Presentation(filepath)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            full_text = " ".join(parts)
            return self._check_and_get_snippet(full_text, keywords, cond)
        except Exception:
            return None

    def search_pdf(self, filepath, keywords, cond):
        try:
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                parts = []
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        parts.append(text)
                full_text = " ".join(parts)
                return self._check_and_get_snippet(full_text, keywords, cond)
        except Exception:
            return None

if __name__ == "__main__":
    root = tk.Tk()
    app = FileSearchApp(root)
    root.mainloop()

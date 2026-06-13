import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

def add_slide(title_text, content_text):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = title_text
    content.text = content_text
    
    # Add a prompt for image
    left = Inches(1.5)
    top = Inches(4.5)
    width = Inches(7)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "【この下にスクリーンショットを貼り付けてください】"
    p.font.size = Pt(20)
    p.font.bold = True
    
    return slide

# Title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "社内メールビューア ver2\n操作マニュアル"
subtitle.text = "※プレースホルダー部分に実際の画面のスクリーンショットを貼り付けて完成させてください。"

# Slide 1
text1 = "1. ルートフォルダ・検索バー（上部）: 検索キーワードや期間の入力、フォルダの指定を行います。\n" \
        "2. フォルダツリー（左側）: サイボウズの「受信箱」や「送信箱」と同じ構成で並びます。\n" \
        "3. メール一覧・本文表示（右側）: 選択したフォルダ内のメールと本文が表示されます。"
add_slide("1. アプリの起動と画面構成", text1)

# Slide 2
text2 = "1. 左側のツリーから、見たいフォルダをクリックします。\n" \
        "2. 右上のリストに、そのフォルダに入っているメールが一覧表示されます。\n" \
        "3. リストの一番上（タイトル、送信者、日時）をクリックすると並べ替えが可能です。\n" \
        "4. 読みたいメールを1回クリックすると、すぐ下に本文が表示されます。"
add_slide("2. メールの閲覧・並べ替え", text2)

# Slide 3
text3 = "1. 画面左上の「検索キーワード」欄に、探したい文字を入力します。\n" \
        "2. カンマ（,）で区切って複数のキーワードを入力できます。\n" \
        "3. 大きな「検索」ボタンをクリックします。\n" \
        "4. 該当するメールが一覧表示され、本文の検索キーワードが黄色く光ります。"
add_slide("3. A. キーワードで探す", text3)

# Slide 4
text4 = "1. 画面の「期間(YYYY/MM/DD)」欄に日付を入力します。\n" \
        "2. 左側に開始日、右側に終了日を入力して「検索」ボタンを押します。\n" \
        "※片方だけの入力でも「〇〇日以降」「〇〇日まで」として検索可能です。"
add_slide("3. B. 日付（期間）で絞り込む", text4)

# Slide 5
text5 = "1. 「検索キーワード」と「期間」に入力されている文字をすべて消して空（カラ）にします。\n" \
        "2. そのまま「検索」ボタンをクリックします。\n" \
        "3. 絞り込みが解除され、すべてのメールが再び表示されます。"
add_slide("3. C. 検索のリセット（全件表示に戻す）", text5)

prs.save('c:\\Users\\フォーレスト026\\MyProject\\tools\\dist\\cybozu_viewer_manual.pptx')
print("PPTX Created Successfully!")

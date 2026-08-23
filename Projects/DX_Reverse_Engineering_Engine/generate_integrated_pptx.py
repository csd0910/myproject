import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
import os
import glob
import re

def extract_section_text(soup, header_keyword):
    for h2 in soup.find_all('h2'):
        if header_keyword in h2.text:
            content = []
            node = h2.find_next_sibling()
            while node and node.name != 'h2':
                if node.name in ['p', 'ul', 'ol', 'div']:
                    text = node.get_text(separator='\n', strip=True)
                    if text:
                        content.append(text)
                node = node.find_next_sibling()
            return "\n".join(content)
    return "情報なし"

def create_integrated_pptx(html_files, output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. タイトルスライド
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    p = tb.text_frame.add_paragraph()
    p.text = "DX業務分析レポート 統合資料"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.name = "Meiryo UI"
    p.alignment = PP_ALIGN.CENTER
    
    sub_p = tb.text_frame.add_paragraph()
    sub_p.text = f"対象日: {', '.join([os.path.basename(f).replace('dx_analysis_report_', '').replace('.html', '') for f in html_files])}"
    sub_p.font.size = Pt(20)
    sub_p.font.name = "Meiryo UI"
    sub_p.alignment = PP_ALIGN.CENTER

    for html_file in html_files:
        with open(html_file, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
        
        date_str = os.path.basename(html_file).replace('dx_analysis_report_', '').replace('.html', '')
        
        # 2. サマリー＆使用ファイルスライド
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        slide2.shapes.title.text = f"[{date_str}] 使用ファイルと手作業のコスト"
        slide2.shapes.title.text_frame.paragraphs[0].font.name = "Meiryo UI"
        
        tb2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        files_text = extract_section_text(soup, "操作した主要ファイル")
        p2 = tb2.text_frame.add_paragraph()
        p2.text = files_text[:1000] + ("..." if len(files_text) > 1000 else "")
        p2.font.name = "Meiryo UI"
        p2.font.size = Pt(14)
        
        # 3. フロー図（テキスト表現）
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])
        slide3.shapes.title.text = f"[{date_str}] 処理フロー図 (Mermaidテキスト)"
        slide3.shapes.title.text_frame.paragraphs[0].font.name = "Meiryo UI"
        
        tb3 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        flow_text = extract_section_text(soup, "現在の処理フロー図")
        p3 = tb3.text_frame.add_paragraph()
        p3.text = flow_text[:1000] + ("..." if len(flow_text) > 1000 else "")
        p3.font.name = "Meiryo UI"
        p3.font.size = Pt(12)

        # 4. 自動化案スライド
        slide4 = prs.slides.add_slide(prs.slide_layouts[5])
        slide4.shapes.title.text = f"[{date_str}] 具体的な自動化案とステップ"
        slide4.shapes.title.text_frame.paragraphs[0].font.name = "Meiryo UI"
        
        tb4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        auto_text = extract_section_text(soup, "具体的な自動化案")
        p4 = tb4.text_frame.add_paragraph()
        p4.text = auto_text[:1000] + ("..." if len(auto_text) > 1000 else "")
        p4.font.name = "Meiryo UI"
        p4.font.size = Pt(14)

    prs.save(output_path)
    print(f"統合PPTXレポートを出力しました: {output_path}")

if __name__ == '__main__':
    # ユーザー指定の2つのHTMLファイル
    html_files = [
        r"C:\AutoAnalysisLogs\daily_reports\dx_analysis_report_20260810.html",
        r"C:\AutoAnalysisLogs\daily_reports\dx_analysis_report_20260811.html"
    ]
    # 実在するファイルのみ対象にする
    valid_files = [f for f in html_files if os.path.exists(f)]
    if valid_files:
        out_path = r"C:\AutoAnalysisLogs\daily_reports\DX_Integrated_Report.pptx"
        create_integrated_pptx(valid_files, out_path)
    else:
        print("指定されたHTMLファイルが見つかりません。")

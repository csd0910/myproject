import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    p = tb.text_frame.add_paragraph()
    p.text = title_text
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.name = "Meiryo UI"
    p.alignment = PP_ALIGN.CENTER
    
    sub_p = tb.text_frame.add_paragraph()
    sub_p.text = subtitle_text
    sub_p.font.size = Pt(20)
    sub_p.font.name = "Meiryo UI"
    sub_p.font.color.rgb = RGBColor(100, 100, 100)
    sub_p.alignment = PP_ALIGN.CENTER

def add_content_slide(title_text, content_list):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title_text
    slide.shapes.title.text_frame.paragraphs[0].font.name = "Meiryo UI"
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item['text']
        p.level = item.get('level', 0)
        p.font.name = "Meiryo UI"
        p.font.size = Pt(20 if p.level == 0 else 16)
        if item.get('bold'): p.font.bold = True

def draw_ipo_flow(slide, title, input_text, process_text, output_text, y_offset=2.0):
    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(y_offset - 0.5), Inches(12), Inches(0.5))
    p = tb.text_frame.add_paragraph()
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.name = "Meiryo UI"
    
    # Input
    shape_in = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_offset + 0.2), Inches(3.5), Inches(2.0))
    shape_in.fill.solid(); shape_in.fill.fore_color.rgb = RGBColor(220, 240, 255)
    shape_in.line.color.rgb = RGBColor(0,0,0)
    tf_in = shape_in.text_frame
    tf_in.text = "【入力 (Input)】\n" + input_text
    tf_in.paragraphs[0].font.bold = True
    for p in tf_in.paragraphs: p.font.name = "Meiryo UI"; p.font.size = Pt(14)
    
    # Arrow 1
    arr1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.1), Inches(y_offset + 0.9), Inches(0.8), Inches(0.6))
    arr1.fill.solid(); arr1.fill.fore_color.rgb = RGBColor(180, 180, 180); arr1.line.color.rgb = RGBColor(180, 180, 180)
    
    # Process
    shape_pr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(y_offset + 0.2), Inches(3.5), Inches(2.0))
    shape_pr.fill.solid(); shape_pr.fill.fore_color.rgb = RGBColor(255, 230, 200)
    shape_pr.line.color.rgb = RGBColor(0,0,0)
    tf_pr = shape_pr.text_frame
    tf_pr.text = "【処理 (Process)】\n" + process_text
    tf_pr.paragraphs[0].font.bold = True
    for p in tf_pr.paragraphs: p.font.name = "Meiryo UI"; p.font.size = Pt(14)
    
    # Arrow 2
    arr2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.6), Inches(y_offset + 0.9), Inches(0.8), Inches(0.6))
    arr2.fill.solid(); arr2.fill.fore_color.rgb = RGBColor(180, 180, 180); arr2.line.color.rgb = RGBColor(180, 180, 180)
    
    # Output
    shape_out = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(y_offset + 0.2), Inches(3.5), Inches(2.0))
    shape_out.fill.solid(); shape_out.fill.fore_color.rgb = RGBColor(220, 255, 220)
    shape_out.line.color.rgb = RGBColor(0,0,0)
    tf_out = shape_out.text_frame
    tf_out.text = "【出力 (Output)】\n" + output_text
    tf_out.paragraphs[0].font.bold = True
    for p in tf_out.paragraphs: p.font.name = "Meiryo UI"; p.font.size = Pt(14)


# Slide 1: Title
add_title_slide("DX Reverse Engineering Engine", "〜 業務の暗黙知を抽出する完全自動化エンジン（ブループリント詳細版） 〜")

# Slide 2: 概要
add_content_slide("システムの概要 (Concept)", [
    {"text": "これは単なるロガーではなく、「業務の暗黙知を抽出する完全自動化エンジン」です。", "bold": True},
    {"text": "人間がやりながら記録し、AIが振り返り、人間が説明を足すことで、業務の完全な姿が浮かび上がります。", "bold": False},
    {"text": "フェーズ構成:", "bold": True},
    {"text": "1. 完全記録（Recording）: バックグラウンドでのPC操作とExcel深層データの抽出", "level": 1},
    {"text": "2. リバース・トレース（Reverse Tracing）: AIによる「結果から逆算した」仮説構築", "level": 1},
    {"text": "3. 作業者への説明要求（Operator Annotation）: セル直結の吹き出しを使ったUI提供", "level": 1},
    {"text": "4. 最終統合と可視化（Final Synthesis）: 客観ログと主観の結合による究極のマニュアル化", "level": 1},
])

# Slide 3: IPO 1
slide_ipo1 = prs.slides.add_slide(prs.slide_layouts[6])
slide_ipo1.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8)).text_frame.text = "フェーズ1・2: ログ抽出とAI推論のフロー (Input -> Process -> Output)"
slide_ipo1.shapes[0].text_frame.paragraphs[0].font.size = Pt(28)
slide_ipo1.shapes[0].text_frame.paragraphs[0].font.name = "Meiryo UI"
slide_ipo1.shapes[0].text_frame.paragraphs[0].font.bold = True

draw_ipo_flow(slide_ipo1, "1. system_logger.py (完全記録フェーズ)",
              "・人間の自然なPC操作\n・Excel上でのコピペ/抽出\n・ソートや色変更",
              "・pynputによるキー打鍵監視\n・win32comによるExcelプロセスへのアタッチ\n・Selection, フィルター, 色情報の強制抽出",
              "・生の操作ログ (system_log_YYYYMMDD.csv)\n※「何を・どこで・どのように」の事実のみ", 1.5)

draw_ipo_flow(slide_ipo1, "2. evaluate_log.py (AIリバース・トレースフェーズ)",
              "・生の操作ログCSV\n・時間的間隔やAlt+Tabの回数\n・対象ファイルのパス遷移",
              "・ログを時系列でチャンク化\n・Gemini APIへ送信\n・「なぜその操作をしたか」の推測・文脈推定",
              "・AI評価済みログ\n(system_log_ai_evaluated_*.csv)\n※「何のために」の仮説が付与されたデータ", 4.5)

# Slide 4: IPO 2
slide_ipo2 = prs.slides.add_slide(prs.slide_layouts[6])
slide_ipo2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8)).text_frame.text = "フェーズ3・4: 図解生成と暗黙知補完のフロー (Input -> Process -> Output)"
slide_ipo2.shapes[0].text_frame.paragraphs[0].font.size = Pt(28)
slide_ipo2.shapes[0].text_frame.paragraphs[0].font.name = "Meiryo UI"
slide_ipo2.shapes[0].text_frame.paragraphs[0].font.bold = True

draw_ipo_flow(slide_ipo2, "3. generate_visual_excel_report.py (説明要求フェーズ)",
              "・操作元のExcelファイル\n・AI評価済みログ (対象セル座標含む)\n・AIによる行動仮説",
              "・win32comで元ファイルを開き別名保存(_v2対応)\n・対象セル(Cell: C3等)を指すCallout図形を生成\n・【作業担当者記載欄】を動的追記",
              "・VisualReport.xlsb (人間が回答を書き込める対話型UI)\n※白背景・黒枠の純正吹き出し付きExcel", 1.5)

draw_ipo_flow(slide_ipo2, "4. 最終統合・可視化フェーズ (Final Synthesis)",
              "・PCの客観的ログ (What/How)\n・VisualReportに作業者が追記した主観コメント (Why)\n・formula_analyzer.pyによる数式構造",
              "・全データを統合解析\n・暗黙のルール(色や目視確認)を明文化\n・自動化可能な要件を抽出",
              "・究極の業務マニュアル 兼 自動化要件定義書\n・(必要に応じて) Python自動化コードの生成", 4.5)

# Slide 5: 関数解析 (Parallel Module)
slide_ipo3 = prs.slides.add_slide(prs.slide_layouts[6])
slide_ipo3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8)).text_frame.text = "技術アーキテクチャ詳細: 関数リバースエンジニアリング"
slide_ipo3.shapes[0].text_frame.paragraphs[0].font.size = Pt(28)
slide_ipo3.shapes[0].text_frame.paragraphs[0].font.name = "Meiryo UI"
slide_ipo3.shapes[0].text_frame.paragraphs[0].font.bold = True

draw_ipo_flow(slide_ipo3, "並行解析モジュール: formula_analyzer.py",
              "・分析対象のExcelファイル\n(複数のシート、数万件のVLOOKUP等)",
              "・全シートの UsedRange.HasFormula を走査\n・数式が含まれるセル座標と実体をメモリ展開\n・操作ログとの突き合わせ",
              "・暗黙の計算ロジックの可視化\n(人間が値を入力した裏で、どの数式が連鎖発火したかの解明)", 2.5)

out_path = r"C:\Users\フォーレスト026\MyProject\DX_Reverse_Engineering_Engine\System_Architecture_Detailed.pptx"
prs.save(out_path)
print(f"PPTX資料を生成しました: {out_path}")
